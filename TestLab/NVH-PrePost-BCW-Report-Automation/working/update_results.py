"""
STEP 7 — NVH Pre/Post BCW Results Updater
Reads dB values from NVH_Report_Input.xlsx (Step 7 columns) and updates
each existing report PPTX with:
  - Actual dB values (replaces XXdB placeholders)
  - Pass / Conditional Fail / Absolute Fail result text
  - Updated Testing Results section text

Run this AFTER you have completed your engineering review and entered
the dB values in the Samples sheet Step-7 columns.
"""

import os
import sys
import glob

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    import openpyxl
except ImportError as e:
    print(f"\nERROR: Missing required library: {e}")
    print("Please ask your administrator to run:  pip install python-pptx openpyxl")
    input("\nPress Enter to close...")
    sys.exit(1)

# ── Paths ────────────────────────────────────────────────────────────
WORKING_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR    = os.path.dirname(WORKING_DIR)
REPORTS_DIR = os.path.join(BASE_DIR, "reports")
INPUT_EXCEL = os.path.join(WORKING_DIR, "NVH_Report_Input.xlsx")

# Slide order matches the 6 dB columns in the Excel Samples sheet
SLIDE_CONDITIONS = [
    "Drive +450Nm (dB)",
    "Drive +75Nm (dB)",
    "Drive -75Nm (dB)",
    "Coast +450Nm (dB)",
    "Coast +75Nm (dB)",
    "Coast -75Nm (dB)",
]

# Slide type labels used in the Testing Results text box
SLIDE_LABELS = ["Drive", "Drive", "Drive", "Coast", "Coast", "Coast"]

# Colours used for the Pass/Fail result text
COLOUR_PASS   = RGBColor(0x00, 0x70, 0xC0)   # blue
COLOUR_CFAIL  = RGBColor(0xFF, 0xC0, 0x00)   # amber
COLOUR_AFAIL  = RGBColor(0xFF, 0x00, 0x00)   # red


# ── Pass/Fail logic ───────────────────────────────────────────────────

def classify(db_value):
    """Return (result_text, colour) for a given dB value."""
    try:
        v = float(db_value)
    except (ValueError, TypeError):
        return None, None
    if v <= 3:
        return "PASS", COLOUR_PASS
    elif v <= 6:
        return "CONDITIONAL FAIL", COLOUR_CFAIL
    else:
        return "ABSOLUTE FAIL", COLOUR_AFAIL


# ── Text helpers ──────────────────────────────────────────────────────

def para_text(para):
    return "".join(r.text for r in para.runs)


def set_para_text(para, new_text, colour=None):
    """Set paragraph text (first run) and optionally change its colour."""
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""
    if colour is not None:
        for r in para.runs:
            r.font.color.rgb = colour


def replace_in_shape(shape, old, new, colour=None):
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = para_text(para)
                if old in full:
                    set_para_text(para, full.replace(old, new), colour)
    except Exception:
        pass
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                replace_in_shape(s, old, new, colour)
    except Exception:
        pass


def replace_in_slide(slide, old, new, colour=None):
    for shape in slide.shapes:
        replace_in_shape(shape, old, new, colour)


# ── Excel reading ─────────────────────────────────────────────────────

def read_samples(wb):
    ws = wb["Samples"]
    headers = [ws.cell(3, j).value for j in range(1, ws.max_column + 1)]
    samples = []
    for row in ws.iter_rows(min_row=4, values_only=True):
        # Skip blank rows and note rows — require both Sample # and Part Number
        if row[0] is None:
            continue
        pn = row[1] if len(row) > 1 else None
        if pn is None or str(pn).strip() == "":
            continue
        sample = {}
        for j, h in enumerate(headers):
            if h:
                val = row[j] if j < len(row) else None
                sample[str(h).strip()] = str(val).strip() if val is not None else ""
        samples.append(sample)
    return samples


# ── Report update ─────────────────────────────────────────────────────

def find_report(pn, sn):
    """Find the report PPTX for a given part number and serial number."""
    pattern = os.path.join(REPORTS_DIR, f"*{pn}*{sn}*.pptx")
    matches = glob.glob(pattern)
    if matches:
        return matches[0]
    # Broader fallback search
    pattern2 = os.path.join(REPORTS_DIR, f"*{sn}*.pptx")
    matches2 = glob.glob(pattern2)
    return matches2[0] if matches2 else None


def update_one_report(report_path, sample):
    prs = Presentation(report_path)

    for slide_idx, condition_key in enumerate(SLIDE_CONDITIONS):
        db_raw = sample.get(condition_key, "")
        if not db_raw:
            continue   # dB not entered yet — skip this slide

        result, colour = classify(db_raw)
        slide_type     = SLIDE_LABELS[slide_idx]
        db_display     = f"{db_raw}dB" if not str(db_raw).endswith("dB") else db_raw

        slide = prs.slides[slide_idx]

        # 1. Replace the large dB display placeholder (XXdB)
        replace_in_slide(slide, "XXdB", db_display, colour)

        # 2. Replace the Testing Results line (Post-100%: XXdB(100%))
        replace_in_slide(slide,
                         f"Post-100%: XXdB(100%)",
                         f"Post-100%: {db_display}(100%)")

        # 3. Add/update the result verdict in the Testing Results box
        #    Pattern: "Drive (Amplitude vs. Target +3dB):" or "Coast ..."
        #    We append the verdict on the same line or update the comments area.
        #    Find any shape that contains the verdict placeholder VERDICT_XX
        replace_in_slide(slide, "VERDICT_XX",
                         result if result else "", colour)

    prs.save(report_path)


# ── Main ──────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  NVH Pre/Post BCW Results Updater — STEP 7")
    print("=" * 60)

    if not os.path.exists(INPUT_EXCEL):
        print(f"\nERROR: Input Excel not found: {INPUT_EXCEL}")
        input("\nPress Enter to close...")
        return

    print("\nReading input file...")
    wb      = openpyxl.load_workbook(INPUT_EXCEL)
    samples = read_samples(wb)

    # Filter to only samples that have at least one dB value filled in
    to_update = [s for s in samples
                 if any(s.get(k, "") for k in SLIDE_CONDITIONS)]

    print(f"  Samples with dB values entered: {len(to_update)}")

    if not to_update:
        print("\n  No dB values found in the Step-7 columns.")
        print("  Please fill in the dB values in the Samples sheet first.")
        input("\nPress Enter to close...")
        return

    print(f"\nUpdating reports in:")
    print(f"  {REPORTS_DIR}\n")

    ok, fail, missing = 0, 0, 0
    for sample in to_update:
        pn = sample.get("Part Number (P/N)", "?")
        sn = sample.get("Serial Number (S/N)", "?")
        report = find_report(pn, sn)

        if not report:
            print(f"  [NOT FOUND]  P/N {pn}  S/N {sn}")
            missing += 1
            continue

        try:
            update_one_report(report, sample)
            print(f"  [UPDATED]  {os.path.basename(report)}")
            ok += 1
        except Exception as e:
            print(f"  [FAIL]  {os.path.basename(report)}  — {e}")
            fail += 1

    print(f"\n{'=' * 60}")
    print(f"  Done.  {ok} updated,  {fail} failed,  {missing} not found.")
    if missing:
        print(f"  Tip: Generate reports first with Step2_Generate_Reports.bat")
    print(f"{'=' * 60}")
    input("\nPress Enter to close...")


if __name__ == "__main__":
    main()
