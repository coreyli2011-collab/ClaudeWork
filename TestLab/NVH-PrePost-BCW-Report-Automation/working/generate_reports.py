"""
STEP 2 — NVH Pre/Post BCW Report Generator
Reads NVH_Report_Input.xlsx and generates one PPTX per sample row.
All test information is filled in automatically. dB values remain as XX placeholders.
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import openpyxl
except ImportError as e:
    print(f"\nERROR: Missing required library: {e}")
    print("Please ask your administrator to run:  pip install python-pptx openpyxl")
    input("\nPress Enter to close...")
    sys.exit(1)

# ── Paths ────────────────────────────────────────────────────────────
WORKING_DIR  = os.path.dirname(os.path.abspath(__file__))
BASE_DIR     = os.path.dirname(WORKING_DIR)
REPORTS_DIR  = os.path.join(BASE_DIR, "reports")
TEMPLATE     = os.path.join(WORKING_DIR, "report_template.pptx")
INPUT_EXCEL  = os.path.join(WORKING_DIR, "NVH_Report_Input.xlsx")


# ── Text replacement helpers ─────────────────────────────────────────

def para_text(para):
    """Return the full text of a paragraph across all its runs."""
    return "".join(r.text for r in para.runs)


def set_para_text(para, new_text):
    """
    Put new_text into the first run and blank out the rest.
    Preserves the formatting (font, size, colour) of the first run.
    """
    if not para.runs:
        return
    para.runs[0].text = new_text
    for r in para.runs[1:]:
        r.text = ""


def replace_in_shape(shape, old, new):
    """Replace old→new in every paragraph of a shape (recursive for groups)."""
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full = para_text(para)
                if old in full:
                    set_para_text(para, full.replace(old, new))
    except Exception:
        pass

    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                replace_in_shape(s, old, new)
    except Exception:
        pass


def replace_in_presentation(prs, replacements):
    """Apply a dictionary of {old: new} replacements to the whole presentation."""
    for slide in prs.slides:
        for shape in slide.shapes:
            for old, new in replacements.items():
                replace_in_shape(shape, old, new)


# ── Excel reading ─────────────────────────────────────────────────────

def read_campaign(wb):
    ws = wb["Campaign_Info"]
    data = {}
    for row in ws.iter_rows(min_row=2, values_only=True):
        if row[0] and row[1] is not None:
            data[str(row[0]).strip()] = str(row[1]).strip()
    return data


def read_samples(wb):
    ws = wb["Samples"]
    # Headers are in row 3
    headers = [ws.cell(3, j).value for j in range(1, ws.max_column + 1)]
    samples = []
    for row in ws.iter_rows(min_row=4, values_only=True):
        # Skip blank rows and note rows — require both Sample # and Part Number
        if row[0] is None:
            continue
        pn = row[1] if len(row) > 1 else None
        if pn is None or str(pn).strip() == "":
            continue
        sample = {}
        for j, h in enumerate(headers):
            if h:
                val = row[j] if j < len(row) else None
                sample[str(h).strip()] = str(val).strip() if val is not None else ""
        samples.append(sample)
    return samples


# ── Report generation ─────────────────────────────────────────────────

def safe_filename(text):
    """Remove characters that are invalid in Windows filenames."""
    for ch in r'\/:*?"<>|':
        text = text.replace(ch, "-")
    return text.strip()


def generate_one_report(campaign, sample):
    prs = Presentation(TEMPLATE)

    pn  = sample.get("Part Number (P/N)", "")
    sn  = sample.get("Serial Number (S/N)", "")
    sno = sample.get("Sample Number", "")
    pub_date = sample.get("Published Date", "")
    pub_by   = campaign.get("Published By", "")

    replacements = {
        # Programme name (header)
        "Stellantis CUSW 2-Speed PV PTU 131mm" : campaign.get("Program Name", ""),

        # Testing Information block
        "NVH Test Order: 282084"                : f"NVH Test Order: {campaign.get('NVH Test Order','')}",
        "NVH Test Dyno: RHTC AWD"               : f"NVH Test Dyno: {campaign.get('Test Dyno','')}",
        "NVH-001: Pre-Test (2/10/2022)"         : f"NVH-001: Pre-Test ({campaign.get('Pre-Test Date','')})",
        "NVH-001: Post-100% (3/2/2022)"         : f"NVH-001: Post-100% ({campaign.get('Post-100% Date','')})",
        "NVH-001: Post-300% (n/a)"              : f"NVH-001: Post-300% ({campaign.get('Post-300% Date','n/a')})",
        "PTL Test Order:  282083"               : f"PTL Test Order:  {campaign.get('PTL Test Order','')}",
        "PTL Sample Number: 0002"               : f"PTL Sample Number: {sno}",

        # Part Information block
        "Test Part Design Level: PV (diff re-sourcing)": f"Test Part Design Level: {campaign.get('Design Level','')}",
        "Test Part P/N: 68333255AE"             : f"Test Part P/N: {pn}",
        "Test Part S/N: T36P22024115524"        : f"Test Part S/N: {sn}",
        "Test Part Ratio: 2.73"                 : f"Test Part Ratio: {campaign.get('Part Ratio','')}",
        "Test Prop Info: PS1247"                : f"Test Prop Info: {campaign.get('Prop Info','')}",

        # Published By (with double-space gap — matches the template exactly)
        "Published By:  C. Li  3/2/2022"        : f"Published By:  {pub_by}  {pub_date}",
    }

    replace_in_presentation(prs, replacements)

    filename = safe_filename(f"NVH_Report_{pn}_SN{sn}.pptx")
    output   = os.path.join(REPORTS_DIR, filename)
    prs.save(output)
    return filename


# ── Main ──────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("  NVH Pre/Post BCW Report Generator — STEP 2")
    print("=" * 60)

    # Check template
    if not os.path.exists(TEMPLATE):
        print(f"\nERROR: Template file not found.")
        print(f"  Expected: {TEMPLATE}")
        print(f"\n  Please prepare your template PPTX (replace actual dB values")
        print(f"  with XXdB), then save it as 'report_template.pptx' in:")
        print(f"  {WORKING_DIR}")
        input("\nPress Enter to close...")
        return

    # Check Excel input
    if not os.path.exists(INPUT_EXCEL):
        print(f"\nERROR: Input Excel not found.")
        print(f"  Expected: {INPUT_EXCEL}")
        print(f"\n  Run 'Create_Input_Template.bat' first to create the Excel file.")
        input("\nPress Enter to close...")
        return

    # Read Excel
    print("\nReading input file...")
    wb       = openpyxl.load_workbook(INPUT_EXCEL)
    campaign = read_campaign(wb)
    samples  = read_samples(wb)

    print(f"  Program  : {campaign.get('Program Name', '(not set)')}")
    print(f"  Samples  : {len(samples)} found")

    if not samples:
        print("\nNo sample rows found in the Samples sheet.")
        input("Press Enter to close...")
        return

    # Generate
    os.makedirs(REPORTS_DIR, exist_ok=True)
    print(f"\nGenerating reports into:")
    print(f"  {REPORTS_DIR}\n")

    ok, fail = 0, 0
    for sample in samples:
        pn = sample.get("Part Number (P/N)", "?")
        sn = sample.get("Serial Number (S/N)", "?")
        try:
            fname = generate_one_report(campaign, sample)
            print(f"  [OK]  {fname}")
            ok += 1
        except Exception as e:
            print(f"  [FAIL] P/N {pn}  S/N {sn}  — {e}")
            fail += 1

    print(f"\n{'=' * 60}")
    print(f"  Done.  {ok} report(s) generated,  {fail} failed.")
    print(f"{'=' * 60}")
    input("\nPress Enter to close...")


if __name__ == "__main__":
    main()
